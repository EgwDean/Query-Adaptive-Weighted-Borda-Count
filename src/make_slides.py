"""Build the lecture-style presentation deck from the study results.

    python src/make_slides.py        # -> slides/query_adaptive_fusion.pptx

Written to be taught from: every concept is defined before it is used, the
maths is spelled out, and the fusion and calibration steps are shown with fully
worked numeric examples. Result numbers are read from the CSVs on disk, so
re-running after a new study refreshes the deck.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from utils import repo_root, load_config, get_paths

INK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x0F, 0x4C, 0x81)
POS = RGBColor(0x2E, 0x7D, 0x32)
NEG = RGBColor(0xC6, 0x28, 0x28)
GREY = RGBColor(0x5F, 0x5F, 0x6F)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
CODEBG = RGBColor(0xEE, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA2, 0x27)

PRIMARY = "score-minmax"
DS_ORDER = ["hotpotqa", "fever", "nfcorpus", "scifact", "msmarco", "fiqa", "quora"]
FLABEL = {"score-minmax": "score fusion", "rrf": "RRF", "borda": "Borda"}


# --------------------------------------------------------------------------- #
# data
# --------------------------------------------------------------------------- #
def load_data(paths):
    rf = paths["router_final"]
    summ = pd.read_csv(os.path.join(rf, "STUDY_SUMMARY.csv"))
    ap = os.path.join(paths["alpha_results"], "oracle_alpha_summary.csv")
    alpha = pd.read_csv(ap) if os.path.exists(ap) else None
    allp = os.path.join(rf, "h2_decision_rule_ALL.csv")
    h2all = pd.read_csv(allp) if os.path.exists(allp) else None
    bva = os.path.join(paths["results"], "figures", "best_vs_mean_alpha.csv")
    bvat = pd.read_csv(bva) if os.path.exists(bva) else None
    return summ, alpha, h2all, bvat


def h2_pooled(h2all):
    m = h2all[h2all.family != "reference"]
    r, c = m[m.rule == "raw"], m[m.rule == "calibrated"]
    sig = "sig_vs_constant" in m.columns
    d = dict(n=len(r), cells=m.groupby(["dataset", "fusion"]).ngroups,
             ds=m.dataset.nunique(),
             raw_beat=int(r.beats_constant.sum()), cal_beat=int(c.beats_constant.sum()))
    if sig:
        d.update(raw_sb=int((r.sig_vs_constant & (r.vs_constant > 0)).sum()),
                 raw_sw=int((r.sig_vs_constant & (r.vs_constant < 0)).sum()),
                 cal_sb=int((c.sig_vs_constant & (c.vs_constant > 0)).sum()),
                 cal_sw=int((c.sig_vs_constant & (c.vs_constant < 0)).sum()))
    return d


def h1_stats(summ):
    from scipy import stats as st
    h = summ[summ.role == "held-out"]
    g = h.groupby("dataset").agg(iqr=("alpha_iqr", "mean"), gain=("gain", "mean"))
    rho, p = st.spearmanr(g.iqr, g.gain)
    r, pp = st.pearsonr(g.iqr, g.gain)
    rc, pc = st.pearsonr(h.alpha_iqr, h.gain)
    return dict(n_ds=len(g), rho=rho, p=p, r=r, pr=pp, cell_r=rc, cell_p=pc,
                n_cells=len(h), table=g.sort_values("iqr", ascending=False))


# --------------------------------------------------------------------------- #
# deck
# --------------------------------------------------------------------------- #
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.part = ""

    def _s(self):
        return self.prs.slides.add_slide(self.blank)

    def _band(self, s):
        b = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.13))
        b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()
        if self.part:
            tb = s.shapes.add_textbox(Inches(10.0), Inches(0.16), Inches(3.2), Inches(0.3))
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = self.part
            r.font.size = Pt(10); r.font.color.rgb = GREY

    def _title(self, s, text, size=27):
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(9.3), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = ACCENT
        r.font.name = "Calibri"

    def note(self, s, t):
        if t:
            s.notes_slide.notes_text_frame.text = t

    def _body(self, s, items, top=1.35, left=0.6, width=12.1, height=5.6, size=17):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for it in items:
            lvl, text = it[0], it[1]
            style = it[2] if len(it) > 2 else ""
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = min(lvl, 4)
            p.space_after = Pt(6 if lvl else 9)
            r = p.add_run(); r.text = text
            r.font.size = Pt(size - 2 * lvl)
            r.font.bold = "b" in style
            r.font.italic = "i" in style
            r.font.name = "Consolas" if "m" in style else "Calibri"
            r.font.color.rgb = (POS if "pos" in style else NEG if "neg" in style
                                else GREY if "muted" in style else INK)
        return tb

    # ---------- slide kinds ----------
    def title_slide(self, title, sub, meta, note=""):
        s = self._s()
        bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.0))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
        p = tf.add_paragraph(); r = p.add_run(); r.text = sub
        r.font.size = Pt(19); r.font.color.rgb = GOLD
        p = tf.add_paragraph(); p.space_before = Pt(18)
        r = p.add_run(); r.text = meta
        r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xBF, 0xC6, 0xD0)
        self.note(s, note); return s

    def part_divider(self, num, title, sub="", note=""):
        self.part = f"Part {num}"
        s = self._s()
        bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = ACCENT; bg.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"Part {num}"
        r.font.size = Pt(17); r.font.color.rgb = RGBColor(0xBF, 0xD8, 0xF0)
        p = tf.add_paragraph(); r = p.add_run(); r.text = title
        r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = WHITE
        if sub:
            p = tf.add_paragraph(); r = p.add_run(); r.text = sub
            r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xDD, 0xE8, 0xF4)
        self.note(s, note); return s

    def bullets(self, title, items, note="", size=17):
        s = self._s(); self._band(s); self._title(s, title)
        self._body(s, items, size=size)
        self.note(s, note); return s

    def math(self, title, formula_lines, items=None, note="", fsize=16):
        """A boxed monospace formula block, optionally followed by bullets."""
        s = self._s(); self._band(s); self._title(s, title)
        h = 0.42 * len(formula_lines) + 0.45
        box = s.shapes.add_shape(1, Inches(0.7), Inches(1.35), Inches(11.9), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = CODEBG
        box.line.color.rgb = RGBColor(0xC5, 0xD0, 0xDE)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.14)
        first = True
        for ln in formula_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = ln
            r.font.size = Pt(fsize); r.font.name = "Consolas"
            r.font.color.rgb = INK
        if items:
            self._body(s, items, top=1.35 + h + 0.25, height=5.4 - h, size=16)
        self.note(s, note); return s

    def table_slide(self, title, header, rows, note="", col_w=None, fontsize=12,
                    top=1.4, bold_col=None, items=None):
        s = self._s(); self._band(s); self._title(s, title)
        nr, nc = len(rows) + 1, len(header)
        th = min(0.34 * nr, 5.3)
        t = s.shapes.add_table(nr, nc, Inches(0.5), Inches(top),
                               Inches(12.3), Inches(th)).table
        if col_w:
            for i, w in enumerate(col_w):
                t.columns[i].width = Inches(w)
        for j, hh in enumerate(header):
            c = t.cell(0, j); c.text = str(hh)
            pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
            run = pr.runs[0]; run.font.size = Pt(fontsize); run.font.bold = True
            run.font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        for i, row in enumerate(rows, 1):
            for j, v in enumerate(row):
                c = t.cell(i, j); c.text = str(v)
                pr = c.text_frame.paragraphs[0]
                pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                run = pr.runs[0]; run.font.size = Pt(fontsize)
                if j == 0 or (bold_col is not None and j == bold_col):
                    run.font.bold = True
                sv = str(v)
                if sv.startswith("+") or sv in ("yes", "GOOD"):
                    run.font.color.rgb = POS
                elif sv.startswith("−") or sv.startswith("-0.") or sv == "no":
                    run.font.color.rgb = NEG
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        if items:
            self._body(s, items, top=top + th + 0.2, height=6.9 - (top + th), size=14)
        self.note(s, note); return s

    def image_slide(self, title, img, note="", caption="", width=10.4, top=1.35):
        s = self._s(); self._band(s); self._title(s, title)
        if img and os.path.exists(img):
            pic = s.shapes.add_picture(img, Inches(0), Inches(top), width=Inches(width))
            pic.left = int((self.prs.slide_width - pic.width) / 2)
        if caption:
            tb = s.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5))
            r = tb.text_frame.paragraphs[0].add_run(); r.text = caption
            r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GREY
        self.note(s, note); return s

    def keypoint(self, kicker, text, note="", color=ACCENT):
        s = self._s(); self._band(s)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.0))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = kicker
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GOLD
        p = tf.add_paragraph(); p.space_before = Pt(10)
        r = p.add_run(); r.text = text
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = color
        self.note(s, note); return s

    def save(self, p):
        self.prs.save(p)


def sgn(x, n=4):
    return f"{x:+.{n}f}"


# --------------------------------------------------------------------------- #
# content
# --------------------------------------------------------------------------- #
def part0(d):
    d.part = ""
    d.title_slide(
        "Query-Adaptive Score Fusion in Hybrid Retrieval",
        "When per-query fusion helps, why it usually does not, and how to make it safe",
        "A complete walkthrough  ·  BEIR  ·  7 datasets × 3 fusion functions",
        "Introduce this as a full walkthrough: we build every concept from scratch, then show "
        "the experiments and the honest results.")
    d.bullets("How to read these slides", [
        (0, "Nothing is assumed. Every term is defined before it is used.", "b"),
        (0, "Formulas appear in grey boxes, always followed by a worked numeric example.", ),
        (1, "If you can follow the example, you can follow the method."),
        (0, "Blue divider slides mark the ten parts.", ),
        (0, "Colour convention used throughout the results:", "b"),
        (1, "green = statistically significant improvement", "pos"),
        (1, "red = statistically significant loss", "neg"),
        (1, "grey = not statistically significant", "muted"),
    ], note="Set expectations. Mediocre-student-proof: definitions first, then maths, then a worked example.")
    d.bullets("Roadmap", [
        (0, "Part 1  — Information retrieval and how we measure it (NDCG)"),
        (0, "Part 2  — The two retrievers: BM25 (lexical) and a dense encoder (semantic)"),
        (0, "Part 3  — Fusion: score fusion, RRF, Borda, with the maths"),
        (0, "Part 4  — The oracle, the α→NDCG curve, and 'headroom'"),
        (0, "Part 5  — The router: features and models"),
        (0, "Part 6  — Calibration: the heart of the contribution, fully worked"),
        (0, "Part 7  — Experimental design and protocol"),
        (0, "Part 8  — The hypotheses we set out to test"),
        (0, "Part 9  — Results"),
        (0, "Part 10 — Correctness, limitations, conclusions"),
    ], note="Give the shape of the talk so nobody is lost.")
    d.keypoint("The whole project in one sentence",
               "Can we pick a better lexical/semantic mix for EACH query\n"
               "than one fixed mix for the whole dataset?",
               "This is the single question. Everything else is machinery for answering it honestly.")


def part1(d):
    d.part_divider(1, "Information retrieval and how we measure it",
                   "Documents, queries, ranking, and the NDCG metric",
                   "We need a shared vocabulary and one metric before anything else.")
    d.bullets("The retrieval task", [
        (0, "We have a corpus: a large collection of documents (here: thousands to 8.8 million).", "b"),
        (0, "A user issues a query — a short piece of text.", "b"),
        (0, "The system must return a ranked list of documents, best first.", "b"),
        (0, "Quality depends almost entirely on what sits at the very top.", "b"),
        (1, "Users look at the first few results and rarely scroll."),
    ], note="Ranking, not classification. Top-of-list quality is what matters.")
    d.bullets("Relevance judgements (qrels)", [
        (0, "To measure quality we need ground truth: which documents are actually relevant.", "b"),
        (0, "A qrel is a triple: (query id, document id, relevance grade).", "m"),
        (0, "Grades are 0 = not relevant, 1 = relevant, sometimes 2 = highly relevant.", ),
        (0, "Only documents judged > 0 count as relevant; unjudged documents are treated as not relevant.", ),
        (0, "BEIR ships qrels split into train / dev / test.", "b"),
        (1, "This split is what lets us train honestly and test once."),
    ], note="Explain qrels plainly: it is the answer key. Note the unjudged=irrelevant assumption.")
    d.bullets("Measuring a ranking: the intuition behind NDCG", [
        (0, "We want a number that rewards putting relevant documents high.", "b"),
        (0, "Three ideas combine:", "b"),
        (1, "1. Gain — a more relevant document is worth more."),
        (1, "2. Discount — a document lower down is worth less."),
        (1, "3. Normalisation — divide by the best possible score, so results land in [0, 1]."),
        (0, "NDCG = Normalised Discounted Cumulative Gain.", "b"),
    ], note="Three ideas: gain, discount, normalise. Say them in order.")
    d.math("NDCG: the formulas",
           ["gain(d)  = 2^rel(d) − 1            rel = relevance grade",
            "discount(i) = 1 / log2(i + 2)      i = 0-based rank position",
            "",
            "DCG@k  = Σ over top-k documents of  gain(d) × discount(i)",
            "IDCG@k = the DCG of the BEST possible ranking",
            "",
            "NDCG@k = DCG@k / IDCG@k            in [0, 1]"],
           [(0, "2^rel − 1 makes a grade-2 document worth 3, a grade-1 worth 1 — highly relevant "
                "documents count much more."),
            (0, "The log discount falls slowly: rank 1 → 1.00, rank 2 → 0.63, rank 3 → 0.50, rank 10 → 0.29.")],
           note="Walk the formula left to right. Stress that IDCG is the perfect ranking, which is why NDCG maxes at 1.")
    d.table_slide("NDCG worked example: k = 5",
                  ["rank i", "document", "rel", "gain = 2^rel−1", "discount 1/log2(i+2)", "contribution"],
                  [["1", "D_a", "1", "1", "1.000", "1.000"],
                   ["2", "D_b", "0", "0", "0.631", "0.000"],
                   ["3", "D_c", "2", "3", "0.500", "1.500"],
                   ["4", "D_d", "0", "0", "0.431", "0.000"],
                   ["5", "D_e", "1", "1", "0.387", "0.387"]],
                  col_w=[1.5, 2.0, 1.2, 2.3, 3.0, 2.3], fontsize=13,
                  items=[(0, "DCG@5 = 1.000 + 0 + 1.500 + 0 + 0.387 = 2.887", "mb"),
                         (0, "Ideal order would be rel = 2, 1, 1 → 3×1.000 + 1×0.631 + 1×0.500 = 4.131", "m"),
                         (0, "NDCG@5 = 2.887 / 4.131 = 0.699", "mb")],
                  note="Do this arithmetic out loud. It demystifies the metric completely.")
    d.bullets("Our metric choices", [
        (0, "Primary metric: NDCG@10.", "b"),
        (1, "Standard in BEIR, and reflects what a user actually sees."),
        (0, "Candidate pool: top_k = 1000 documents per retriever.", "b", ),
        (1, "Retrieve deep, evaluate shallow — a good document can be rescued from rank 800 by fusion."),
        (0, "We also report MRR@100 and Recall@100 as secondary numbers.", ),
        (1, "MRR = 1 / rank of the first relevant document."),
        (1, "Recall@100 = fraction of all relevant documents found in the top 100."),
    ], note="eval_k=10 is the decision metric; top_k=1000 is the working pool. Do not confuse them.")


def part2(d):
    d.part_divider(2, "The two retrievers",
                   "BM25 (lexical) and all-mpnet-base-v2 (semantic)",
                   "Two fundamentally different ways to decide what is relevant.")
    d.bullets("Two families of retriever", [
        (0, "Lexical (sparse): match the actual words.", "b"),
        (1, "Query 'Einstein 1921 Nobel Prize' → find documents containing those terms."),
        (1, "Strong on names, IDs, rare terms, exact phrasing."),
        (1, "Fails when the document says the same thing in different words."),
        (0, "Semantic (dense): match the meaning.", "b"),
        (1, "Query 'how do I fix a leaky tap' matches a document about 'repairing a dripping faucet'."),
        (1, "Strong on paraphrase and intent."),
        (1, "Fails on rare exact tokens it never learned, e.g. product codes."),
        (0, "They fail in different places — that is what makes combining them worthwhile.", "b", "pos"),
    ], note="This complementarity is the entire premise of hybrid retrieval.")
    d.bullets("BM25: the idea", [
        (0, "A bag-of-words scoring function; the standard lexical baseline for decades.", "b"),
        (0, "Three intuitions:", "b"),
        (1, "1. A query term appearing in a document is evidence of relevance."),
        (1, "2. A RARE term is stronger evidence than a common one ('quantum' beats 'the')."),
        (1, "3. Repeating a term helps, but with diminishing returns — 20 occurrences is not 20× better."),
        (1, "4. Long documents contain more words by chance, so they must be penalised."),
    ], note="Give the four intuitions before showing the formula; the formula is just these written down.")
    d.math("BM25: the formula",
           ["                            f(t,D) × (k1 + 1)",
            "BM25(D,Q) = Σ  IDF(t) × ──────────────────────────────",
            "            t∈Q          f(t,D) + k1 × (1 − b + b×|D|/avgdl)",
            "",
            "f(t,D) = how often term t occurs in document D",
            "|D|    = document length,   avgdl = average document length",
            "IDF(t) = rarity weight of term t (high for rare terms)"],
           [(0, "k1 controls term-frequency saturation — how fast repeats stop helping.", ),
            (0, "b controls length normalisation: b = 0 ignores length, b = 1 fully normalises.", ),
            (0, "IDF is what makes rare terms count more.", )],
           note="Map each intuition onto its symbol: IDF=rarity, k1=saturation, b=length.")
    d.math("BM25 settings we used",
           ["method       = lucene",
            "k1           = 0.8      (default is 1.2 — we tuned lower)",
            "b            = 0.4      (default is 0.75 — less length penalty)",
            "use_stemming = true     ('running' and 'runs' → 'run')"],
           [(0, "Tuned by grid search on hotpotqa: k1 ∈ {0.4 … 2.0} × b ∈ {0.3 … 1.0} × stemming.", ),
            (0, "Limitation we disclose: these values were then INHERITED by every other dataset, "
                "not re-tuned per dataset.", "muted")],
           note="Be upfront that BM25 params were tuned once and inherited. It is a real limitation.")
    d.bullets("Dense retrieval: the idea", [
        (0, "A neural encoder maps any text to a vector (an 'embedding') of 768 numbers.", "b"),
        (0, "Training pushes texts with similar MEANING to nearby vectors.", "b"),
        (0, "We embed every document once, offline, and store the vectors.", ),
        (0, "At query time we embed the query and find the nearest document vectors.", ),
        (0, "Model: sentence-transformers/all-mpnet-base-v2", "mb"),
        (1, "768 dimensions, max 384 word-pieces per text, vectors L2-normalised."),
    ], note="Emphasise the offline/online split: embedding the corpus is the expensive one-time cost.")
    d.math("Dense scoring: cosine similarity",
           ["                    u · v",
            "cos(u, v) = ───────────────────      ∈ [−1, 1]",
            "                 ‖u‖ × ‖v‖",
            "",
            "If all vectors are L2-normalised (‖u‖ = ‖v‖ = 1):",
            "",
            "cos(u, v) = u · v        a plain dot product"],
           [(0, "We normalise at encoding time, so scoring is just a matrix multiply — very fast on a GPU.", ),
            (0, "Higher cosine = more semantically similar.", )],
           note="Normalising turns cosine into a dot product, which is why the GPU search is a single matmul.")
    d.bullets("Scale: what this actually costs", [
        (0, "msmarco has 8,841,823 documents. Embedding them all is the expensive step.", "b"),
        (0, "8.8M × 768 dimensions × 4 bytes (float32) ≈ 27 GB.", "m"),
        (0, "We stored msmarco in float16 instead → ≈ 13.5 GB.", "m", "pos"),
        (0, "Written to disk as a memory-mapped array, so it never has to fit in RAM.", ),
        (0, "Similarity is computed in chunks of 50,000 documents to bound GPU memory.", ),
    ], note="Explain why engineering choices (fp16, memmap, chunking) exist: an 8.8M corpus does not fit anywhere.")


def part3(d):
    d.part_divider(3, "Fusion", "Combining two ranked lists into one",
                   "Now we have two lists per query. How do we merge them?")
    d.bullets("The fusion problem", [
        (0, "For one query we now have TWO ranked lists of 1000 documents each.", "b"),
        (0, "We must produce ONE ranking to show the user.", "b"),
        (0, "Two broad approaches:", "b"),
        (1, "Score fusion — combine the numeric scores."),
        (1, "Rank fusion — ignore scores, combine only the positions."),
        (0, "In both cases a weight α decides how much each retriever counts.", "b"),
    ], note="Frame the choice: use the numbers, or only the order.")
    d.bullets("Problem: the two score scales are incompatible", [
        (0, "BM25 scores are unbounded positive numbers — maybe 8.3, 6.1, 5.9…", "b"),
        (0, "Cosine similarities live in [−1, 1] — maybe 0.71, 0.68, 0.66…", "b"),
        (0, "Adding them directly would let BM25 dominate purely because its numbers are bigger.", "b", "neg"),
        (0, "Fix: normalise each retriever's scores per query, onto a common [0, 1] range.", "b", "pos"),
    ], note="This is why normalisation exists. Without it the weight alpha would be meaningless.")
    d.math("Min-max normalisation (applied per query, per retriever)",
           ["             s − min(s)",
            "norm(s) = ─────────────────        result ∈ [0, 1]",
            "           max(s) − min(s)",
            "",
            "best document in the list  → 1.0",
            "worst document in the list → 0.0",
            "if all scores are equal    → 0.0 for everything"],
           [(0, "Applied to each retriever's own top-1000, independently, for every query.", ),
            (0, "A document retrieved by only one retriever gets 0 from the other side.", )],
           note="Per query and per retriever. That locality is important — scales differ across queries too.")
    d.math("Score fusion: the convex combination",
           ["fuse(d) = α × norm(bm25_score(d))  +  (1 − α) × norm(dense_score(d))",
            "",
            "α = 1.0   →  pure lexical  (BM25 only)",
            "α = 0.0   →  pure semantic (dense only)",
            "α = 0.5   →  equal mixture",
            "",
            "The weights α and (1−α) sum to 1  →  'convex' combination"],
           [(0, "α is THE quantity this whole project is about.", "b"),
            (0, "Normally α is one fixed number for the entire dataset. We ask whether it should "
                "change per query.", "b")],
           note="Land this hard: alpha is the object of study.")
    d.table_slide("Score fusion worked example (α = 0.6)",
                  ["doc", "BM25 raw", "norm BM25", "dense raw", "norm dense", "0.6×b + 0.4×d", "rank"],
                  [["D1", "8.3", "1.00", "0.55", "0.20", "0.68", "1"],
                   ["D2", "6.1", "0.44", "0.71", "1.00", "0.66", "2"],
                   ["D3", "5.9", "0.38", "0.66", "0.70", "0.51", "3"],
                   ["D4", "4.3", "0.00", "0.51", "0.00", "0.00", "4"]],
                  col_w=[1.1, 1.8, 1.9, 1.8, 1.9, 2.6, 1.2], fontsize=13,
                  items=[(0, "BM25 min=4.3 max=8.3 → D1 = (8.3−4.3)/(8.3−4.3) = 1.00, D2 = (6.1−4.3)/4.0 = 0.44", "m"),
                         (0, "D1 wins on lexical, D2 wins on semantic. At α = 0.6 lexical counts more, so D1 edges ahead.", ),
                         (0, "At α = 0.3 the order would flip to D2, D1 — the weight genuinely changes the answer.", "b")],
                  note="Point out the flip: alpha is not cosmetic, it reorders results.")
    d.math("Rank fusion 1: Reciprocal Rank Fusion (RRF)",
           ["                    1                         1",
            "RRF(d) = α × ───────────────  +  (1−α) × ───────────────",
            "              K + rank_b(d)                K + rank_d(d)",
            "",
            "K = 60 (standard constant),  rank is 1-based",
            "a document missing from a list contributes 0 from that side"],
           [(0, "Only the POSITION matters. The score that produced the position is discarded.", "b"),
            (0, "K = 60 damps the top: rank 1 → 1/61, rank 2 → 1/62. Very flat.", )],
           note="RRF is popular because it is scale-free and robust. But flatness is also its weakness.")
    d.math("Rank fusion 2: Borda count",
           ["Borda(d) = α × (N − rank_b(d))  +  (1−α) × (N − rank_d(d))",
            "",
            "N = 1000 (list length), rank is 0-based",
            "",
            "top of a list      → N − 0   = 1000 points",
            "bottom of a list   → N − 999 = 1 point",
            "absent from a list → 0 points"],
           [(0, "Linear in rank, unlike RRF's reciprocal.", ),
            (0, "Absent scores 0, which is worse than being last — that is the intended penalty.", )],
           note="Borda spreads points linearly; RRF concentrates them at the top. Different shapes.")
    d.table_slide("The same query under all three fusions (α = 0.5)",
                  ["doc", "BM25 rank", "dense rank", "score fusion", "RRF", "Borda"],
                  [["D1", "1", "4", "0.60", "0.0157", "997.5"],
                   ["D2", "2", "1", "0.72", "0.0163", "998.5"],
                   ["D3", "3", "2", "0.54", "0.0160", "997.5"],
                   ["D4", "4", "3", "0.00", "0.0158", "996.5"]],
                  col_w=[1.2, 2.2, 2.2, 2.4, 2.2, 2.1], fontsize=13,
                  items=[(0, "Score fusion sees that D1 is FAR ahead on BM25 and D4 is junk; the rank "
                             "fusions only see 1st, 2nd, 3rd, 4th.", "b"),
                         (0, "RRF's numbers are nearly identical to each other — it has thrown away how "
                             "confident each retriever was.", "neg")],
                  note="This table is the argument for score fusion in one picture.")
    d.keypoint("Why score fusion is our primary method",
               "Rank fusion discards score MAGNITUDE —\n"
               "the information saying 'both of these are excellent'\n"
               "rather than merely 'this one is first'.",
               "Credit where due: this is settled in the literature (Bruch, Gai & Ingber, TOIS 2023). "
               "We cite it and do NOT claim it as our contribution.")
    d.bullets("Our position on the three fusion functions", [
        (0, "Score fusion (min-max) is the PRIMARY method throughout.", "b", "pos"),
        (0, "RRF and Borda are included only as standard BASELINES.", "b"),
        (0, "That score fusion beats them is established by Bruch, Gai & Ingber (ACM TOIS 2023) — "
            "we cite it, we do not claim it.", "muted"),
        (0, "So why run all three at all?", "b"),
        (1, "To test whether our CONCLUSIONS depend on the fusion function (this becomes H3)."),
        (1, "A finding that only holds for one fusion rule would be much weaker."),
    ], note="Pre-empt 'why bother with RRF/Borda' — the answer is robustness testing, not competition.")


def part4(d, bvat):
    d.part_divider(4, "The oracle and the headroom",
                   "How good could a perfect α possibly be?",
                   "Before building anything, find out whether there is anything to win.")
    d.bullets("The normal way: one global α", [
        (0, "Standard practice: pick ONE α and use it for every query in the dataset.", "b"),
        (0, "How: try every α on a tuning split and keep whichever maximises average NDCG@10.", ),
        (0, "We sweep α over a grid of 101 values: 0.00, 0.01, 0.02, …, 1.00.", "m"),
        (0, "We call the winner α* — the best constant. It is our BASELINE.", "b", "pos"),
        (0, "Beating α = 0.5 proves nothing. Beating α* is the real test.", "b"),
    ], note="Insist on this: the honest baseline is a properly TUNED constant, not a naive 0.5.")
    d.bullets("The oracle: the best α for each individual query", [
        (0, "Now do the same sweep, but per query, and let each query keep its own best α.", "b"),
        (0, "That per-query best is the ORACLE α.", "b"),
        (0, "It is a cheat: it uses the answer key to choose α.", "b", "neg"),
        (0, "So it is not a method — it is a CEILING. No real router can beat it.", "b"),
        (0, "Its value is diagnostic: it tells us how much per-query routing could possibly win.", ),
    ], note="Students must not think the oracle is a system. It is an upper bound computed with the labels.")
    d.math("The α→NDCG curve (computed once per query)",
           ["for each α in {0.00, 0.01, …, 1.00}:",
            "     curve[α] = NDCG@10 of the ranking produced at that α",
            "",
            "oracle_α        = argmax over α of curve[α]",
            "oracle_ndcg     = max curve[α]",
            "alpha_sensitivity = max(curve) − min(curve)"],
           [(0, "Storing the WHOLE curve (not just the best point) is a key engineering decision.", "b"),
            (0, "Any α a model later predicts can then be scored by a table lookup — no re-ranking.", ),
            (0, "This is what made a 21-cell study computationally feasible.", "pos")],
           note="The stored curve is why everything downstream is fast: prediction -> lookup, not re-retrieval.")
    d.table_slide("One query's α→NDCG curve (abridged to 5 of 101 points)",
                  ["α", "0.00", "0.25", "0.50", "0.75", "1.00"],
                  [["NDCG@10", "0.30", "0.45", "0.60", "0.50", "0.35"]],
                  col_w=[2.4, 1.9, 1.9, 1.9, 1.9, 1.9], fontsize=14,
                  items=[(0, "oracle_α = 0.50  (the argmax)", "mb"),
                         (0, "oracle_ndcg = 0.60", "m"),
                         (0, "alpha_sensitivity = 0.60 − 0.30 = 0.30", "m"),
                         (0, "A FLAT curve would mean α barely matters for this query — its oracle label "
                             "is then almost arbitrary, so we down-weight it during training.", )],
                  note="Explain alpha_sensitivity as a sample weight: flat curve = unreliable label.")
    d.bullets("Complementarity, measured", [
        (0, "Collect the oracle α of every query and look at the SPREAD of that distribution.", "b"),
        (0, "We summarise spread by the interquartile range (IQR = 75th − 25th percentile).", "m"),
        (0, "Wide spread → different queries genuinely want different mixes → routing has something to do.", "b", "pos"),
        (0, "Zero spread → every query wants the same mix → a constant already captures everything.", "b", "neg"),
        (0, "This single number becomes the x-axis of our first hypothesis.", "b"),
    ], note="IQR of oracle alpha = our operational definition of retriever complementarity.")
    d.math("Headroom: the size of the prize",
           ["static_best = NDCG@10 using the single best constant α*",
            "oracle      = mean over queries of each query's own best NDCG@10",
            "",
            "headroom = oracle − static_best",
            "",
            "gain     = router − static_best        (what we actually achieved)",
            "% headroom captured = gain / headroom × 100"],
           [(0, "Headroom is the absolute maximum any router could add.", "b"),
            (0, "If headroom is small, even a perfect router cannot impress.", )],
           note="Define headroom carefully — the results section keeps referring to it.")
    d.keypoint("A trap worth an entire slide",
               "The best constant α is NOT the average of the per-query oracle αs.",
               "This is the mean-alpha fallacy. It is the mechanism behind our main result, so spend time here.")
    d.math("Why the average is the wrong answer",
           ["argmax of the MEAN curve   ≠   MEAN of the per-query argmaxes",
            "",
            "Most queries want a LOW α (they prefer the dense retriever).",
            "A minority want a HIGH α — and they gain a LOT from getting it.",
            "",
            "Averaging the αs is dragged down by the majority.",
            "Maximising average NDCG is pulled UP by the minority's large gains."],
           [(0, "This is a Jensen-style effect: the objective is non-linear in α.", ),
            (0, "It matters because a model trained on oracle-α labels learns the conditional MEAN — "
                "exactly the wrong target.", "b", "neg")],
           note="This is the intellectual core. The regressor targets E[alpha|x]; the optimum is elsewhere.")
    if bvat is not None and len(bvat):
        rows = []
        for ds in DS_ORDER:
            r = bvat[(bvat.dataset == ds) & (bvat.fusion == PRIMARY)]
            if len(r):
                r = r.iloc[0]
                rows.append([ds, f"{r.alpha_mean:.3f}", f"{r.alpha_star:.2f}",
                             f"{r.alpha_star - r.alpha_mean:+.3f}", f"{r.penalty_mean:.4f}"])
        d.table_slide("Measured on our data: mean α vs the α that actually wins",
                      ["dataset", "mean oracle α", "best α*", "difference", "NDCG lost if you use the mean"],
                      rows, col_w=[2.6, 2.6, 2.0, 2.2, 3.5], fontsize=13,
                      items=[(0, "α* is higher than the mean oracle α in almost every case.", "b"),
                             (0, "Using the mean costs up to 0.126 NDCG — vastly more than our router's "
                                 "entire gain of about 0.006.", "b", "neg")],
                      note="Quantifies the trap. The penalty dwarfs every other effect in the study.")


def part5(d):
    d.part_divider(5, "The router", "Predicting α from things we can see at query time",
                   "Now we try to actually predict a good alpha, without the answer key.")
    d.bullets("What the router has to do", [
        (0, "Input: only what is available at query time — the query text and the two ranked lists.", "b"),
        (0, "It may NOT see relevance judgements. Those exist only for training.", "b", "neg"),
        (0, "Output: an α for this query.", "b"),
        (0, "Hard constraint: it must be CHEAP. It runs on every query, before results are shown.", "b"),
        (1, "Our final routers cost roughly one microsecond per query."),
    ], note="The inference-time constraint drives many later decisions, including which models we allow.")
    d.bullets("Where the features come from", [
        (0, "Key insight: the SHAPE of a retriever's score distribution reveals its confidence.", "b", "pos"),
        (0, "A retriever that found one clear winner looks different from one returning mush.", ),
        (1, "Confident: top score far above the rest, high variance, low entropy."),
        (1, "Unsure: all top scores nearly equal, low variance, high entropy."),
        (0, "If BM25 looks confident and dense looks unsure, we should lean lexical — and vice versa.", "b"),
        (0, "This field is called Query Performance Prediction (QPP).", "muted"),
    ], note="Motivate the features conceptually before listing them. Confidence has a measurable shape.")
    d.table_slide("Feature family 1 — score distribution (computed for BOTH retrievers)",
                  ["feature", "meaning", "high value means"],
                  [["top_score", "the single best score", "strong top match"],
                   ["margin", "score[1] − score[2]", "a clear winner"],
                   ["norm_margin", "margin ÷ top score", "clear winner, scale-free"],
                   ["sigma_k", "std-dev of the top-100 scores", "scores well spread out"],
                   ["entropy", "entropy of softmaxed scores", "flat, undifferentiated list"],
                   ["smv", "score magnitude variance", "uneven score mass"],
                   ["robust_sigma", "std-dev after trimming outliers", "spread, outlier-proof"],
                   ["zscore_top", "(top − mean) ÷ std", "top stands out from the pack"],
                   ["zscore_margin", "margin ÷ std", "gap is large relative to noise"]],
                  col_w=[2.7, 5.0, 4.6], fontsize=12,
                  note="These are the workhorses. Every one is a cheap statistic of numbers we already have.")
    d.table_slide("Feature families 2 and 3 — coherence and agreement",
                  ["feature", "meaning", "high value means"],
                  [["autocorr", "score/embedding autocorrelation in the top-W", "results are mutually consistent"],
                   ["apair_ratio", "top-W vs bottom-W embedding similarity", "top results form a tight cluster"],
                   ["query_centroid_cos", "cosine between query and corpus centroid", "generic, non-specific query"],
                   ["jaccard", "overlap of the two top-1000 lists", "the retrievers agree on candidates"],
                   ["kendall_tau", "rank correlation on the shared documents", "they also agree on the ORDER"],
                   ["ql", "query length in words", "longer, more specific query"]],
                  col_w=[3.0, 5.2, 4.1], fontsize=12,
                  items=[(0, "Difference features: d_entropy, d_smv, d_sigma_k, d_zscore_top, d_zscore_margin", "mb"),
                         (0, "Each is simply (BM25 value − dense value) — an explicit 'who looks more "
                             "confident?' signal.", )],
                  note="Agreement features ask whether the two retrievers even disagree. Difference features compare confidence directly.")
    d.bullets("31 features in total", [
        (0, "13 score-distribution features × 2 retrievers = 26", "m"),
        (0, "+ 5 difference features (BM25 − dense)", "m"),
        (0, "+ jaccard, kendall_tau, ql, query_centroid_cos", "m"),
        (0, "All computable from what we already have — no extra retrieval, no extra model calls.", "b", "pos"),
        (0, "We later prune these 31 down to just 3 (Part 5, ablation).", "b"),
    ], note="31 candidates, later cut to 3. The pruning is a result in itself.")
    d.bullets("How we frame the learning problem", [
        (0, "The label is the oracle α — a number in [0, 1]. Three ways to learn it:", "b"),
        (0, "regression — predict α directly as a continuous value.", "mb"),
        (0, "binary — predict P(α > 0.5): 'does this query prefer lexical?'", "mb"),
        (0, "multibin — chop [0,1] into 11 bins, predict which bin, take the expected α.", "mb"),
        (0, "We test ALL THREE for every model family, because the best framing is not obvious.", "b"),
        (0, "Important: the framing is about the LABEL. It is separate from the decision rule (Part 6).", "b", "neg"),
    ], note="Students confuse framing with decision rule constantly. Flag it here and again in Part 6.")
    d.table_slide("Model families screened",
                  ["family", "type", "included?"],
                  [["LightGBM", "gradient-boosted trees", "yes"],
                   ["XGBoost", "gradient-boosted trees", "yes"],
                   ["CatBoost", "gradient-boosted trees", "yes (re-screen)"],
                   ["HistGradientBoosting", "gradient-boosted trees", "yes"],
                   ["RandomForest", "bagged trees", "yes"],
                   ["ExtraTrees", "bagged trees", "yes"],
                   ["ElasticNet", "linear (regression only)", "yes"],
                   ["LogisticRegression", "linear (classification)", "yes"],
                   ["MLP", "small neural network", "yes (re-screen)"],
                   ["SVM", "kernel method", "NO — O(n²) training"],
                   ["k-NN", "instance based", "NO — O(n) at INFERENCE"]],
                  col_w=[3.4, 4.4, 4.5], fontsize=12,
                  items=[(0, "SVM and k-NN are excluded for cost, not accuracy: a router that is slow at "
                             "inference defeats the entire purpose.", "b")],
                  note="Justify exclusions on inference cost. That is a design principle, not laziness.")
    d.math("Hyperparameter search: Optuna with TPE",
           ["for each (family, framing) pair:",
            "     run an INDEPENDENT Optuna study of 30 trials",
            "     objective = mean NDCG@10 on the dev split",
            "",
            "TPE = Tree-structured Parzen Estimator (Bayesian search):",
            "  it models which regions of the space produced good scores",
            "  and samples more often from those regions"],
           [(0, "Independent studies mean every family gets the same budget — a fair best-vs-best "
                "comparison, not a lucky-default comparison.", "b"),
            (0, "Selection happens on dev. Test is untouched.", "pos")],
           note="Fairness argument: equal trial budget per family.")
    d.math("Feature selection: greedy backward elimination",
           ["start with all 31 features",
            "repeat:",
            "     for each feature f still in the set:",
            "          score the model WITHOUT f",
            "     permanently drop whichever f hurt the least",
            "until 3 features remain",
            "",
            "then pick the SMALLEST set statistically tied with the best point"],
           [(0, "Why backward and not 'rank by importance'? Redundancy.", "b"),
            (0, "Two near-duplicate features each look useless alone (the other covers for it), so "
                "importance ranking would wrongly discard both.", "neg"),
            (0, "Ties are broken toward dropping the more EXPENSIVE feature.", )],
           note="The redundancy argument is the key justification for backward elimination.")
    d.bullets("Result of the selection: three features per fusion", [
        (0, "score fusion → margin_bm25, entropy_bm25, smv_dense", "mb"),
        (0, "RRF → ql, smv_dense, d_entropy", "mb"),
        (0, "Borda → ql, sigma_k_dense, d_entropy", "mb"),
        (0, "From 31 candidates down to 3, with no significant loss in quality.", "b", "pos"),
        (0, "Honest caveat: these are 3 specs chosen once on hotpotqa and reused everywhere — "
            "NOT 21 independent confirmations.", "b", "neg"),
    ], note="Pre-empt over-reading the feature-stability table. Three specs, reused. Say it plainly.")


def part6(d):
    d.part_divider(6, "Calibration — the decision layer",
                   "Turning a model's output into a fusion weight",
                   "This is the heart of the contribution. Go slowly.")
    d.keypoint("The central problem",
               "A trained model's output is NOT a fusion weight.",
               "Everything in this part follows from this one observation.")
    d.bullets("Why the raw output fails", [
        (0, "A regressor trained on oracle-α labels learns the conditional MEAN of α.", "b"),
        (0, "But Part 4 showed the NDCG-optimal α is far ABOVE that mean.", "b", "neg"),
        (0, "A classifier's probability is on a probability scale, not an α scale, at all.", "b"),
        (0, "In practice model outputs are SQUASHED into a narrow band around their average.", "b"),
        (1, "The ORDER of the predictions is often informative…"),
        (1, "…but the SCALE is wrong."),
        (0, "Using such a number directly as α systematically under-weights the lexical retriever.", "b", "neg"),
    ], note="Two distinct failures: wrong target (the mean) and wrong scale (squashed). Both fixed by calibration.")
    d.bullets("The fix: histogram binning", [
        (0, "Use the model ONLY to rank queries. Throw its actual numbers away.", "b", "pos"),
        (0, "Three steps:", "b"),
        (1, "1. Sort queries by the model's output and cut them into equal-sized bins."),
        (1, "2. For each bin, average the α→NDCG curves of the queries in it."),
        (1, "3. That bin emits the α maximising its OWN averaged curve."),
        (0, "The bin→α table is learned on a held-out slice of the training data.", "b"),
        (0, "At inference: predict → find bin → emit that bin's stored α. Pure lookup.", "b"),
    ], note="Three steps. Repeat them: rank, average curves, take each bin's argmax.")
    # ---- worked example ----
    d.bullets("Worked example — the setup", [
        (0, "To keep the arithmetic doable we shrink everything:", "b"),
        (0, "9 queries (real: thousands)", "m"),
        (0, "α grid of 5 points: 0, 0.25, 0.5, 0.75, 1.0 (real: 101 points)", "m"),
        (0, "3 bins (real: 10, 20 or 50)", "m"),
        (0, "Each query has an α→NDCG curve, as defined in Part 4.", ),
    ], note="Say clearly this is a shrunk toy so the numbers can be checked by hand.")
    d.table_slide("Worked example — the 9 queries' curves",
                  ["query", "α=0", "α=0.25", "α=0.5", "α=0.75", "α=1.0", "its own best α"],
                  [["q1", "0.60", "0.55", "0.40", "0.30", "0.20", "0.00"],
                   ["q2", "0.50", "0.48", "0.35", "0.25", "0.15", "0.00"],
                   ["q3", "0.70", "0.65", "0.50", "0.40", "0.30", "0.00"],
                   ["q4", "0.30", "0.45", "0.60", "0.50", "0.35", "0.50"],
                   ["q5", "0.25", "0.40", "0.55", "0.45", "0.30", "0.50"],
                   ["q6", "0.35", "0.50", "0.58", "0.52", "0.40", "0.50"],
                   ["q7", "0.20", "0.30", "0.45", "0.55", "0.65", "1.00"],
                   ["q8", "0.15", "0.25", "0.40", "0.50", "0.60", "1.00"],
                   ["q9", "0.25", "0.35", "0.48", "0.58", "0.68", "1.00"]],
                  col_w=[1.4, 1.6, 1.8, 1.7, 1.8, 1.7, 2.3], fontsize=12,
                  note="q1-q3 want dense, q4-q6 want a mix, q7-q9 want lexical. A perfect router would separate them.")
    d.table_slide("Worked example — step 0: the constant baseline",
                  ["α", "0.00", "0.25", "0.50", "0.75", "1.00"],
                  [["mean NDCG over all 9", "0.3667", "0.4367", "0.4789", "0.4500", "0.4033"]],
                  col_w=[3.6, 1.8, 1.8, 1.8, 1.8, 1.8], fontsize=13,
                  items=[(0, "Best constant α* = 0.50, giving mean NDCG = 0.4789", "mb"),
                         (0, "Every query gets α = 0.50. This is the number any router must beat.", "b")],
                  note="Compute the baseline first so the comparison later is meaningful.")
    d.table_slide("Worked example — step 1: the model's raw outputs",
                  ["query", "q1", "q2", "q3", "q4", "q5", "q6", "q7", "q8", "q9"],
                  [["raw output", "0.42", "0.44", "0.46", "0.48", "0.50", "0.52", "0.54", "0.56", "0.58"]],
                  col_w=[2.2, 1.12, 1.12, 1.12, 1.12, 1.12, 1.12, 1.12, 1.12, 1.12], fontsize=12,
                  items=[(0, "The ORDER is perfect: q1 (most dense-preferring) → q9 (most lexical-preferring).", "b", "pos"),
                         (0, "But every value is squashed into [0.42, 0.58] — nothing near the 0.0 and 1.0 "
                             "that queries actually want.", "b", "neg")],
                  note="Good ordering, useless scale. This is exactly what real models produce.")
    d.bullets("Worked example — step 2: what the RAW rule scores", [
        (0, "Use each raw output directly as α. All values round to the α = 0.50 grid point.", "b"),
        (0, "So every query is scored at α = 0.50 — identical to the constant baseline.", "b"),
        (0, "Raw rule mean NDCG = 0.4789. Gain over the constant = 0.0000.", "mb", "neg"),
        (0, "The model's genuine knowledge of the ORDER has been completely wasted.", "b", "neg"),
    ], note="Raw ties the constant here; in the real experiments it often does actively worse.")
    d.table_slide("Worked example — step 3: bin the queries, average their curves",
                  ["bin", "queries", "α=0", "α=0.25", "α=0.5", "α=0.75", "α=1.0", "bin's best α"],
                  [["0 (lowest)", "q1,q2,q3", "0.600", "0.560", "0.417", "0.317", "0.217", "0.00"],
                   ["1 (middle)", "q4,q5,q6", "0.300", "0.450", "0.577", "0.490", "0.350", "0.50"],
                   ["2 (highest)", "q7,q8,q9", "0.200", "0.300", "0.443", "0.543", "0.643", "1.00"]],
                  col_w=[1.9, 2.0, 1.3, 1.5, 1.4, 1.5, 1.4, 1.6], fontsize=12,
                  items=[(0, "Bin 0 average at α=0: (0.60 + 0.50 + 0.70) / 3 = 0.600  ← the row maximum", "m"),
                         (0, "Each bin takes the argmax of its OWN averaged curve — not the average of its αs.", "b")],
                  note="Show one average being computed by hand, then say the rest follow identically.")
    d.math("Worked example — the learned calibration table",
           ["bin 0  →  α = 0.00",
            "bin 1  →  α = 0.50",
            "bin 2  →  α = 1.00"],
           [(0, "This tiny table is the entire learned decision layer. It is what gets frozen to disk.", "b"),
            (0, "Note that it has recovered the full α range 0.0 … 1.0 from outputs that only ever "
                "spanned 0.42 … 0.58.", "b", "pos")],
           note="The table is the artefact. Calibration stretched a compressed signal back onto the real axis.")
    d.table_slide("Worked example — step 4: the scoreboard",
                  ["rule", "how α is chosen", "mean NDCG@10", "vs constant"],
                  [["best constant α*", "0.50 for every query", "0.4789", "—"],
                   ["RAW output", "the model's number, as-is", "0.4789", "+0.0000"],
                   ["CALIBRATED", "the bin's learned α", "0.6067", "+0.1278"]],
                  col_w=[2.9, 4.3, 2.6, 2.4], fontsize=14,
                  items=[(0, "Bin 0 at α=0 → 0.60, 0.50, 0.70; bin 1 at α=0.5 → 0.60, 0.55, 0.58; "
                             "bin 2 at α=1 → 0.65, 0.60, 0.68", "m"),
                         (0, "Mean = 5.46 / 9 = 0.6067", "mb"),
                         (0, "Same model. Same predictions. Only the decision rule changed.", "b", "pos")],
                  note="Land the punchline: identical model, identical predictions, radically different result.")
    d.keypoint("The safety property",
               "If the model has NO signal, every bin picks the same α,\n"
               "and calibration reduces exactly to the constant baseline.",
               "This is why the method is safe to deploy: its worst case is 'no change', by construction.")
    d.bullets("Why the safety property holds", [
        (0, "Suppose the model's ordering is pure noise.", "b"),
        (0, "Then each bin is a random subset of queries.", ),
        (0, "So each bin's averaged curve ≈ the global averaged curve.", ),
        (0, "So each bin's argmax = the global argmax = α*.", ),
        (0, "Every bin emits α*, which IS the constant baseline.", "b", "pos"),
        (0, "Caveat we verified empirically: this guarantee holds on the data the bins are FITTED on. "
            "On unseen queries a small loss is still possible.", "b", "neg"),
    ], note="Give the argument, then the honest caveat. In-sample guarantee, out-of-sample it can slip slightly.")
    d.math("Guard rail: minimum queries per bin",
           ["n_bins = min(requested_bins,  n_queries // 50)",
            "",
            "MIN_QUERIES_PER_BIN = 50"],
           [(0, "With too few queries per bin, each bin's 'best α' is fitted to noise.", "b", "neg"),
            (0, "Real failure we hit: scifact with ~8 queries per bin scored −0.044 — a significant LOSS.", "neg"),
            (0, "After the floor: −0.003, not significant. The method degrades toward the constant "
                "instead of breaking.", "pos")],
           note="A real bug turned into a designed safety mechanism, with before/after numbers.")
    d.bullets("Framing vs decision rule — do not confuse them", [
        (0, "FRAMING = how the label is built and what the model predicts.", "b"),
        (1, "regression → a number; binary → one probability; multibin → a bin distribution."),
        (0, "DECISION RULE = what we do with that prediction.", "b"),
        (1, "raw → use it directly as α;  calibrated → bin and look up."),
        (0, "They are independent choices. Every framing can be calibrated.", "b", "pos"),
        (0, "Example: our Borda router is logreg|binary AND calibrated with 20 bins.", "m"),
    ], note="Recurring student confusion. The binary/multibin label says nothing about calibration.")


def part7(d, summ):
    d.part_divider(7, "Experimental design", "Datasets, splits, and the honesty protocol",
                   "How the experiment was set up so the results can be believed.")
    d.table_slide("The seven datasets",
                  ["dataset", "domain", "corpus", "train", "dev", "test", "role"],
                  [["hotpotqa", "multi-hop Wikipedia QA", "5.2M", "85,000", "5,447", "7,405", "DEVELOPMENT"],
                   ["fever", "fact verification", "5.4M", "109,810", "6,666", "6,666", "held-out"],
                   ["msmarco", "web search", "8.8M", "502,939", "6,980", "43", "held-out"],
                   ["quora", "duplicate questions", "523K", "—", "5,000", "10,000", "held-out"],
                   ["fiqa", "financial QA", "57K", "5,500", "500", "648", "held-out"],
                   ["nfcorpus", "medical IR", "3.6K", "2,590", "324", "323", "held-out"],
                   ["scifact", "scientific claims", "5K", "809", "—", "300", "held-out"]],
                  col_w=[2.0, 3.4, 1.5, 1.7, 1.4, 1.4, 2.2], fontsize=12,
                  items=[(0, "Chosen to SPAN the complementarity axis — from wide oracle-α spread "
                             "(hotpotqa) to exactly zero (quora).", "b"),
                         (0, "quora is the deliberate NEGATIVE CONTROL: we predict no gain there.", "b")],
                  note="Spanning the x-axis is deliberate. quora is designed to fail, which makes success elsewhere meaningful.")
    d.bullets("Why these seven, and not others", [
        (0, "Every held-out dataset must have a non-test split to refit weights and calibration.", "b"),
        (0, "That rules out the many BEIR datasets shipped test-only:", ),
        (1, "arguana, scidocs, trec-covid, dbpedia-entity, webis-touche2020, climate-fever, nq…"),
        (0, "Of BEIR's datasets with a usable train split, we use ALL of them.", "b", "pos"),
        (0, "msmarco needed special handling: BEIR's test split is only 43 queries, so we evaluate on "
            "dev (6,980), which is standard community practice.", "b"),
    ], note="Pre-empt 'why not more datasets': we used everything that qualifies. That is a real constraint.")
    d.math("Split discipline — the rule we never break",
           ["TRAIN  →  fit model weights and the calibration table",
            "DEV    →  select everything (family, hyperparameters, features)",
            "TEST   →  opened exactly ONCE, at the very end, for the final number",
            "",
            "No decision anywhere in this project was made by looking at test."],
           [(0, "If you tune on test, your reported number is optimistic and meaningless.", "b", "neg"),
            (0, "hotpotqa is the DEVELOPMENT dataset — the one place design decisions were made, "
                "which we disclose.", "b")],
           note="The credibility of every number rests on this slide.")
    d.bullets("Frozen-spec inheritance: the strongest part of the protocol", [
        (0, "All model/feature selection happened on hotpotqa ONLY (sections 5, 6, 7).", "b"),
        (0, "Every held-out dataset SKIPS those sections entirely.", "b"),
        (0, "Each inherits hotpotqa's frozen spec: family, framing, hyperparameters, feature set.", "b"),
        (0, "Only two things are refit per dataset:", "b"),
        (1, "the model weights (on that dataset's own train split)"),
        (1, "the calibration table (on that dataset's own α→NDCG curves)"),
        (0, "So nothing is ever SELECTED on held-out data — only fitted.", "b", "pos"),
    ], note="This is what makes the held-out results a genuine test of the DESIGN, not just of the weights.")
    d.bullets("A subtlety: datasets with only two splits", [
        (0, "scifact has no dev split. quora has no train split. msmarco's dev is used for evaluation.", "b"),
        (0, "Naively the code then fits and evaluates on the SAME queries.", "b", "neg"),
        (0, "We found this caused a real distortion: raw outputs looked far better than they are, "
            "because a fine-grained output can memorise the fitting queries.", "b", "neg"),
        (0, "Fix: deterministically carve the available split into two disjoint halves.", "b", "pos"),
        (0, "After the fix, raw's apparent wins on those datasets largely disappeared.", "b"),
    ], note="Show that we hunted for our own leakage and found some. That builds trust.")
    d.table_slide("The pipeline: ten sections, one command",
                  ["#", "section", "what it does"],
                  [["0", "download", "fetch the BEIR dataset"],
                   ["1", "embed", "encode the whole corpus to vectors (the expensive step)"],
                   ["2", "tune_bm25", "grid-search k1 / b / stemming (off by default)"],
                   ["3", "retrieve", "top-1000 from BOTH retrievers, with raw scores, cached"],
                   ["4", "dataset", "compute 31 features + the α→NDCG curve + oracle α"],
                   ["5", "screen", "model families × framings (Optuna, on dev)"],
                   ["6", "ablate", "greedy backward feature elimination"],
                   ["7", "rescreen", "families × framings × feature-set sizes"],
                   ["8", "final_fit", "refit on full train, then FREEZE the router"],
                   ["9", "benchmark", "evaluate on test — once"]],
                  col_w=[0.8, 2.4, 9.1], fontsize=12,
                  items=[(0, "Sections 5–7 run ONLY on hotpotqa. Held-out datasets run 0–4, 8, 9.", "b"),
                         (0, "Every section skips if its outputs exist → the whole study is resumable.", )],
                  note="Sections 5-7 are the selection sections. That is exactly what held-out datasets skip.")
    d.math("Key configuration values",
           ["top_k              = 1000     candidate pool per retriever",
            "eval_k             = 10       PRIMARY METRIC = NDCG@10",
            "alpha grid         = 0.00 … 1.00 step 0.01   (101 points)",
            "features.window    = 100      top-W scores used for distribution features",
            "router.train_subset= 10000    queries used to fit the router",
            "router.n_trials    = 30       Optuna trials per (family, framing)",
            "n_calib_bins       ∈ {10, 20, 50}",
            "calib_fraction     = 0.2      share of train held out for calibration",
            "bootstrap_resamples= 1000     for all confidence intervals",
            "seed               = 42       everything is deterministic"],
           note="Everything needed to reproduce. Point at the seed: the study is fully deterministic.")
    d.bullets("Measuring significance: the paired bootstrap", [
        (0, "Two systems are scored on the SAME queries, so compare them query by query.", "b"),
        (0, "Method:", "b"),
        (1, "1. Compute d = NDCG(router) − NDCG(baseline) for every query."),
        (1, "2. Resample the queries with replacement 1,000 times; take the mean of d each time."),
        (1, "3. Report the 2.5th and 97.5th percentiles — a 95% confidence interval."),
        (0, "Significant = that interval excludes zero.", "b", "pos"),
        (0, "Pairing cancels per-query difficulty, so it detects wins that separate CIs would miss.", ),
    ], note="Explain WHY paired: differencing the same query removes the dominant source of variance.")
    d.math("Paired bootstrap — a tiny example",
           ["router   = [0.60, 0.50, 0.70, 0.55, 0.65]",
            "baseline = [0.55, 0.52, 0.60, 0.55, 0.58]",
            "d        = [+0.05, −0.02, +0.10, 0.00, +0.07]     mean = +0.040",
            "",
            "resample d 1000× → 95% CI = [+0.006, +0.074]",
            "the interval excludes 0  →  SIGNIFICANT"],
           [(0, "If the CI had been [−0.01, +0.09], the same mean would NOT be significant.", ),
            (0, "Small test sets give wide intervals — which is why our 300-query datasets rarely "
                "reach significance.", "b")],
           note="Connects directly to why nfcorpus/scifact/fiqa come out grey in the results.")


def part8(d):
    d.part_divider(8, "The hypotheses", "What we set out to test, and why it changed",
                   "An honest account of how the research question evolved.")
    d.bullets("Where this started", [
        (0, "The original thesis idea: build a query-adaptive hybrid retriever and show it WINS.", "b"),
        (0, "A router predicts α per query; we expected to beat the fixed-α baseline.", ),
        (0, "That is a classic 'our system is better' paper.", "muted"),
    ], note="Start with the naive framing so the pivot has impact.")
    d.bullets("What we actually found", [
        (0, "Against a properly TUNED constant α, the adaptive gains were small.", "b", "neg"),
        (0, "On some datasets they were zero, or slightly negative.", "b", "neg"),
        (0, "A weak 'we win by 0.006' paper would be unconvincing and probably not general.", "muted"),
        (0, "So we changed the question.", "b", "pos"),
    ], note="Do not hide this. The reframe is the intellectually honest move and makes a better paper.")
    d.keypoint("The reframe",
               "From 'our system wins'\nto 'WHEN does this help, WHY does it usually not,\nand HOW do we make it safe?'",
               "A conditional, falsifiable study is worth more than a marginal system claim.")
    d.bullets("H1 — the scaling hypothesis", [
        (0, "Claim: the gain from per-query α grows with retriever COMPLEMENTARITY.", "b"),
        (0, "Complementarity is measured as the IQR of the oracle-α distribution.", "m"),
        (0, "Prediction: wide spread → real gains; zero spread → no gain at all.", "b"),
        (0, "This is falsifiable — quora (IQR = 0.00) must show no gain, or H1 is wrong.", "b"),
    ], note="State the falsifier explicitly. That is what makes it a hypothesis and not a slogan.")
    d.bullets("H2 — the decision-layer hypothesis", [
        (0, "Claim: a router's RAW output is not a fusion weight; used directly it loses to a constant.", "b"),
        (0, "Claim: histogram-binning calibration fixes this, and by construction cannot do worse.", "b"),
        (0, "Prediction: across many model families, raw underperforms and calibrated does not.", "b"),
        (0, "This is the methodological contribution.", "b", "pos"),
    ], note="H2 is the novel part. H1 and H3 are context that make H2 meaningful.")
    d.bullets("H3 — the invariance hypothesis", [
        (0, "Claim: the H1/H2 patterns are not artefacts of one fusion function.", "b"),
        (0, "Prediction: the same behaviour appears under score fusion, RRF, and Borda.", "b"),
        (0, "This is a robustness check: a finding true for only one fusion rule is fragile.", ),
    ], note="H3 is the cheapest hypothesis to test since we already compute all three fusions.")
    d.math("The experimental matrix",
           ["7 datasets  ×  3 fusion functions  =  21 cells",
            "",
            "each cell:  build features → fit router → calibrate → benchmark on test",
            "",
            "plus the H2 experiment:",
            "18 router configurations × 2 decision rules × 7 datasets = 252 runs"],
           note="Give the scale. This is a lot of compute, run unattended over several days.")


def part9(d, summ, h2all, figs):
    d.part_divider(9, "Results", "What the 21 cells and 252 H2 runs actually showed",
                   "Now the evidence, hypothesis by hypothesis.")
    # H2
    d.bullets("H2 first — it is the strongest result", [
        (0, "We screened 18 router configurations (6 families × 3 framings) under BOTH rules…", "b"),
        (0, "…on every one of the 7 datasets, with paired-bootstrap confidence intervals.", "b"),
        (0, "Question: does the raw output beat a properly tuned constant α?", "b"),
    ], note="Set up the question before the numbers.")
    if h2all is not None:
        p = h2_pooled(h2all)
        rows = [["raw output used as α", f"{p.get('raw_sb','?')}/{p['n']}", f"{p.get('raw_sw','?')}/{p['n']}"],
                ["calibrated", f"{p.get('cal_sb','?')}/{p['n']}", f"{p.get('cal_sw','?')}/{p['n']}"]]
        d.table_slide(f"H2 pooled over {p['ds']} datasets ({p['n']} configurations)",
                      ["decision rule", "significantly BETTER than constant", "significantly WORSE"],
                      rows, col_w=[4.0, 4.6, 3.7], fontsize=15, top=1.6,
                      items=[(0, "Raw is roughly 18× more likely to significantly HURT than to help.", "b", "neg"),
                             (0, "Calibration is far more likely to help than hurt, and never by much when it slips.", "b", "pos"),
                             (0, "Honest wording: say calibration 'almost never' loses, not 'never'.", "b")],
                      note="The headline table. Note we corrected 'never' to 'almost never' after finding one significant loss.")
    d.image_slide("H2 on the development dataset, configuration by configuration",
                  figs.get("h2"), width=11.0,
                  caption="Each x position is one model configuration. Red = raw, green = calibrated, dashed = the constant.",
                  note="Every red point below the line, every green point above. Same models, only the decision rule differs.")
    d.bullets("H2 — the mechanism, confirmed on all 7 datasets", [
        (0, "Recall the mean-α fallacy from Part 4.", "b"),
        (0, "We recorded each raw router's MEAN prediction and compared it to the NDCG-optimal α.", "b"),
        (0, "In 7 of 7 datasets the raw predictions sit BELOW the optimum. Mean gap ≈ −0.19.", "b", "neg"),
        (1, "scifact: predicted 0.256 vs optimal 0.630 — a gap of −0.374."),
        (1, "fever: predicted 0.354 vs optimal 0.590."),
        (0, "So we do not merely observe that raw fails — we show WHY it fails.", "b", "pos"),
    ], note="Mechanism beats observation. This turns H2 from an empirical curiosity into an explanation.")
    d.image_slide("The mean-α fallacy, measured", figs.get("bva"), width=12.2,
                  caption="Left: every cell sits above the diagonal. Right: the NDCG cost of using the mean.",
                  note="Left panel is the proof; right panel is the cost. This is the mechanism slide.")
    # H1
    hs = h1_stats(summ)
    d.image_slide("H1: gain versus complementarity", figs.get("h1"), width=10.6,
                  caption="Each point is one (dataset, fusion) cell. X marks the development dataset, excluded from the fit.",
                  note="Upward trend. fever top-right, quora bottom-left in red.")
    rows = [[ds, f"{r.iqr:.3f}", sgn(r.gain)] for ds, r in hs["table"].iterrows()]
    d.table_slide("H1 at the dataset level (the statistically honest view)",
                  ["dataset", "oracle-α IQR", "mean gain"], rows,
                  col_w=[3.2, 3.0, 3.0], fontsize=13, top=1.5,
                  items=[(0, f"Spearman ρ = {hs['rho']:+.3f}, p = {hs['p']:.4f}  (n = {hs['n_ds']} datasets)", "mb", "pos"),
                         (0, f"Pearson r = {hs['r']:+.3f}, p = {hs['pr']:.3f} — the RANK relationship is "
                             f"much clearer than the linear one.", "m")],
                  note="Lead with Spearman. The relationship is monotonic, not linear.")
    d.bullets("An important statistical caveat we must state ourselves", [
        (0, f"Treating all {hs['n_cells']} cells as independent gives r = {hs['cell_r']:+.3f}, "
            f"p = {hs['cell_p']:.4f} — which looks great.", "b"),
        (0, "But that is PSEUDO-REPLICATION.", "b", "neg"),
        (1, "The 3 fusion cells of a dataset share the same queries, retrievers and judgements."),
        (1, "They are not 3 independent pieces of evidence."),
        (0, f"The honest unit is the dataset: n = {hs['n_ds']}, not {hs['n_cells']}.", "b"),
        (0, "We report the dataset-level statistic as the headline.", "b", "pos"),
    ], note="Reviewers WILL catch this. Catching it ourselves is far better.")
    d.image_slide("H1 holds inside each fusion function separately", figs.get("h1f"), width=12.4,
                  caption="Fitted independently per fusion — the trend is not an artefact of pooling.",
                  note="Supports H1 and H3 simultaneously.")
    d.bullets("H1 — the two-factor refinement", [
        (0, "Complementarity is NECESSARY but not SUFFICIENT.", "b"),
        (0, "You also need enough data to learn and to detect the effect:", "b"),
        (1, "fever — IQR 0.48, 6,666 test queries → +0.031, significant.", "pos"),
        (1, "nfcorpus — IQR 0.49, only 323 test queries → +0.003, not significant.", "muted"),
        (0, "msmarco is the decisive case: 6,980 queries (plenty of power) but LOW spread → "
            "gain only +0.0027.", "b"),
        (0, "So the ceiling is set by complementarity, not by sample size.", "b", "pos"),
    ], note="msmarco settles the confound: with power but no complementarity, the gain stays tiny.")
    # H3
    d.image_slide("H3: the same pattern under all three fusion functions", figs.get("gbf"), width=11.8,
                  caption="★ = significant. Error bars are 95% paired-bootstrap intervals.",
                  note="Wide error bars on the small datasets are the visual proof that they are underpowered, not negative.")
    # main table
    prim = summ[summ.fusion == PRIMARY].set_index("dataset")
    rows = []
    for ds in DS_ORDER:
        if ds in prim.index:
            r = prim.loc[ds]
            rows.append([ds, f"{r.alpha_iqr:.2f}", f"{int(r.n_queries):,}", f"{r.static_best:.4f}",
                         f"{r.router:.4f}", f"{r.oracle:.4f}", sgn(r.gain),
                         "yes" if r.significant else "no"])
    d.table_slide("Main results — score fusion (test split)",
                  ["dataset", "IQR", "queries", "constant α*", "router", "oracle", "gain", "sig?"],
                  rows, col_w=[1.9, 1.1, 1.5, 1.9, 1.6, 1.6, 1.6, 1.0], fontsize=12,
                  items=[(0, "Significant gains only where complementarity is high — exactly H1's prediction.", "b"),
                         (0, "quora shows a significant NEGATIVE gain: with nothing to route on, routing "
                             "slightly hurts.", "b", "neg")],
                  note="The money table. Small numbers, honestly reported, and consistent with the hypothesis.")
    d.bullets("The negative control worked", [
        (0, "quora: oracle-α IQR = 0.00. Every query wants the same mix.", "b"),
        (0, "H1 predicts no gain. Observed: a small but SIGNIFICANT loss (−0.003) in all three fusions.", "b", "neg"),
        (0, "With 10,000 test queries there is ample power, so this is a real effect, not noise.", ),
        (0, "A negative control that behaves as predicted makes the positive results far more credible.", "b", "pos"),
        (0, "Caveat: −0.003 is statistically real but practically negligible. Report effect sizes, "
            "not just stars.", "muted"),
    ], note="Explain why a predicted failure strengthens the paper.")
    d.bullets("Safety: the calibration guard rail in action", [
        (0, "scifact, score fusion, before the fix: −0.044 — a significant LOSS.", "b", "neg"),
        (0, "Diagnosis: 20 bins over ~160 queries ≈ 8 queries per bin. Each bin's α was fitted to noise.", "b"),
        (0, "Fix: require ≥ 50 queries per bin, reducing the bin count when data is scarce.", "b"),
        (0, "After: −0.003, not significant. The method degrades gracefully toward the constant.", "b", "pos"),
        (0, "We report this as a deployment requirement, not as a footnote.", ),
    ], note="A failure diagnosed, fixed, and measured. That is what makes it a contribution.")


def part10(d):
    d.part_divider(10, "Correctness, limitations, conclusions",
                   "What we checked, what we cannot claim, and what it all means",
                   "Finish with the caveats. Credibility comes from the limitations slide.")
    d.bullets("We audited our own code and found two real bugs", [
        (0, "Bug 1 — phantom documents in rank fusion.", "b", "neg"),
        (1, "Queries made only of stop-words produce an empty BM25 result."),
        (1, "The placeholder row was harmless for score fusion but gave FAKE rank points under Borda/RRF."),
        (1, "Fixed: a document counts only if the retriever actually scored it."),
        (0, "Bug 2 — evaluation leakage on two-split datasets.", "b", "neg"),
        (1, "quora and scifact fitted and evaluated on the same queries."),
        (1, "This inflated the RAW rule specifically: 20 of 21 apparent 'raw wins' came from these two."),
        (1, "Fixed: carve deterministic disjoint halves."),
        (0, "Both were found by us, before publication, and both are reported.", "b", "pos"),
    ], note="Do not hide bugs. Finding and fixing them is evidence of rigour.")
    d.bullets("Limitations — stated up front", [
        (0, "BM25 hyperparameters were tuned on hotpotqa and inherited, not re-tuned per dataset.", "muted"),
        (0, "H1 rests on 6 independent held-out datasets. That is few for a scaling claim.", "muted"),
        (0, "We report per-cell confidence intervals with no multiple-comparison correction.", "muted"),
        (0, "The three small datasets (300–650 test queries) are underpowered by construction.", "muted"),
        (0, "Feature 'stability' is 3 specs reused, not 21 independent selections.", "muted"),
        (0, "One dense encoder and one lexical retriever; other pairings may differ.", "muted"),
    ], note="Every one of these is something a reviewer would otherwise raise. Own them.")
    d.bullets("What we would do with more compute", [
        (0, "Add more independent datasets — the only real way to strengthen H1.", "b"),
        (1, "Note: adding more FUSIONS does not help; it inflates n without adding evidence."),
        (0, "Apply FDR or Bonferroni correction across the 21 cells.", ),
        (0, "Test a second dense encoder to check the findings are not model-specific.", ),
        (0, "Sweep bin count against dataset size to turn the safety floor into a systematic result.", ),
    ], note="Show we know exactly what the next experiment is.")
    d.bullets("Contributions", [
        (0, "1. A decision layer that makes router-predicted fusion weights safe, replicated across "
            "7 datasets with confidence intervals.", "b", "pos"),
        (0, "2. The mechanism behind it: models trained on oracle-α labels target the conditional mean, "
            "which is systematically below the NDCG-optimal α (7/7 datasets).", "b", "pos"),
        (0, "3. A condition for when query-adaptive fusion is worth doing: retriever complementarity, "
            "with data volume as a second requirement.", "b", "pos"),
        (0, "4. A validated safety mechanism, from a diagnosed real failure.", "b", "pos"),
    ], note="Four contributions, each with evidence in this deck.")
    d.bullets("Practical guidance for a practitioner", [
        (0, "Before building any router, measure the oracle-α IQR on your own data.", "b"),
        (1, "Near zero? Use a tuned constant α and spend your time elsewhere."),
        (0, "Always tune the constant baseline properly. Beating α = 0.5 proves nothing.", "b"),
        (0, "Never feed a model's raw output in as a fusion weight. Calibrate it.", "b"),
        (0, "Never pick a constant by averaging oracle αs — maximise NDCG instead.", "b"),
        (0, "Keep at least ~50 queries per calibration bin.", "b"),
    ], note="The takeaway slide for anyone who has to build this in production.")
    d.keypoint("Summary",
               "Adaptive fusion helps only where the retrievers genuinely disagree.\n"
               "A router's output must be calibrated, never used raw.\n"
               "The pattern holds across all three fusion functions.",
               "Three sentences. If the audience remembers only this slide, that is enough.")
    d.part = ""
    d.title_slide("Thank you", "Questions welcome",
                  "Code, configuration, and all result tables are in the repository.",
                  "Open the floor. Likely questions: dataset count, multiple comparisons, other encoders.")


# --------------------------------------------------------------------------- #
def build():
    root = repo_root()
    cfg = load_config()
    paths = get_paths(cfg)
    summ, alpha, h2all, bvat = load_data(paths)

    outdir = os.path.join(root, "slides")
    figdir = os.path.join(outdir, "figures")
    os.makedirs(figdir, exist_ok=True)

    import plot_study as PS
    made = PS.build_all(paths, outdir=figdir)
    figs = dict(h2=made.get("h2"), h1=made.get("h1_combined"),
                h1f=made.get("h1_by_fusion"), gbf=made.get("gain_by_fusion"),
                bva=made.get("best_vs_mean_alpha"))

    d = Deck()
    part0(d)
    part1(d)
    part2(d)
    part3(d)
    part4(d, bvat)
    part5(d)
    part6(d)
    part7(d, summ)
    part8(d)
    part9(d, summ, h2all, figs)
    part10(d)

    out = os.path.join(outdir, "query_adaptive_fusion.pptx")
    d.save(out)
    print(f"[slides] wrote {out}  ({len(list(d.prs.slides))} slides)")
    return out


if __name__ == "__main__":
    build()
